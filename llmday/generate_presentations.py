#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Generator prezentacji PowerPoint z plików SUMMARY markdown
Tworzy żywe, interesujące prezentacje pełne wiedzy z LLM Day
"""

import sys
import platform
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# UTF-8 handling for Windows
if platform.system() == 'Windows':
    try:
        sys.stdout.reconfigure(encoding='utf-8')
    except:
        pass

class PresentationGenerator:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)

    def add_title_slide(self, title, subtitle):
        """Dodaj slajd tytułowy"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[0])
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]

        title_shape.text = title
        subtitle_shape.text = subtitle

        # Styling
        title_shape.text_frame.paragraphs[0].font.size = Pt(44)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    def add_section_header(self, title):
        """Dodaj slajd nagłówka sekcji"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])  # Blank

        left = Inches(1)
        top = Inches(3)
        width = Inches(8)
        height = Inches(1.5)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.text = title

        # Styling
        p = text_frame.paragraphs[0]
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 102, 0)
        p.alignment = PP_ALIGN.CENTER

    def add_content_slide(self, title, content_items):
        """Dodaj slajd z treścią (lista punktów)"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])

        title_shape = slide.shapes.title
        body_shape = slide.placeholders[1]

        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(36)
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

        text_frame = body_shape.text_frame
        text_frame.clear()

        for item in content_items:
            p = text_frame.add_paragraph()
            p.text = item
            p.level = 0
            p.font.size = Pt(20)
            p.space_before = Pt(6)

    def add_two_column_slide(self, title, left_content, right_content):
        """Dodaj slajd z dwiema kolumnami"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])  # Blank

        # Title
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(9)
        height = Inches(1)

        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(36)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

        # Left column
        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(4.5), Inches(5))
        left_frame = left_box.text_frame
        left_frame.word_wrap = True

        for item in left_content:
            p = left_frame.add_paragraph()
            p.text = item
            p.font.size = Pt(18)
            p.space_before = Pt(6)

        # Right column
        right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.8), Inches(4.5), Inches(5))
        right_frame = right_box.text_frame
        right_frame.word_wrap = True

        for item in right_content:
            p = right_frame.add_paragraph()
            p.text = item
            p.font.size = Pt(18)
            p.space_before = Pt(6)

    def add_quote_slide(self, quote, attribution=""):
        """Dodaj slajd z cytatem"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])  # Blank

        left = Inches(1.5)
        top = Inches(2.5)
        width = Inches(7)
        height = Inches(3)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True

        # Quote
        p = text_frame.paragraphs[0]
        p.text = f'"{quote}"'
        p.font.size = Pt(32)
        p.font.italic = True
        p.font.color.rgb = RGBColor(51, 51, 51)
        p.alignment = PP_ALIGN.CENTER

        # Attribution
        if attribution:
            p2 = text_frame.add_paragraph()
            p2.text = f"— {attribution}"
            p2.font.size = Pt(24)
            p2.font.color.rgb = RGBColor(102, 102, 102)
            p2.alignment = PP_ALIGN.CENTER
            p2.space_before = Pt(12)

    def save(self, filename):
        """Zapisz prezentację"""
        self.prs.save(filename)
        print(f"[SUCCESS] Saved: {filename}")


def generate_intel_presentation():
    """Generuj prezentację Intel OpenVino"""
    gen = PresentationGenerator()

    # Slide 1: Title
    gen.add_title_slide(
        "Intel OpenVino",
        "Lokalne AI bez Clouda\nLLM Day 2024"
    )

    # Slide 2: Main Thesis
    gen.add_section_header("Nie potrzebujesz clouda!")

    # Slide 3: Key Numbers
    gen.add_content_slide(
        "Kluczowe Liczby",
        [
            "10 tokens/s = szybkość czytania człowieka",
            "NPU zużywa 60% energii GPU (40% oszczędności!)",
            "87% RAM można dedykować dla AI (BIOS setting)",
            "20W NPU vs 26-27W GPU",
            "Apache 2.0 - commercial use OK!"
        ]
    )

    # Slide 4: Live Demo
    gen.add_content_slide(
        "Live Demo - BEZ INTERNETU!",
        [
            "Laptop Intel Core Ultra + GPU Arc 140V",
            "Model: Qwen3 8B",
            "Rezultat: 10 tokens/s streaming",
            "Task Manager pokazuje NPU w akcji",
            "Banana classification demo - działa!",
            "Real-time responses"
        ]
    )

    # Slide 5: Quote
    gen.add_quote_slide(
        "Pytanie jest, czy musisz szybciej? W większości przypadków, to jest absolutnie wystarczająco."
    )

    # Slide 6: Cloud vs Local
    gen.add_two_column_slide(
        "Cloud vs Local (OpenVino)",
        [
            "CLOUD:",
            "• Wymaga internetu",
            "• Dane wysyłane na zewnątrz",
            "• Ciągłe koszty API",
            "• Network lag",
            "• Ograniczona kontrola"
        ],
        [
            "LOCAL (OpenVino):",
            "• Offline capable",
            "• Dane zostają lokalnie",
            "• Jednorazowy koszt HW",
            "• Szybkie (local)",
            "• Pełna kontrola",
            "• 60% energii (NPU!)"
        ]
    )

    # Slide 7: Use Cases
    gen.add_content_slide(
        "Praktyczne Zastosowania",
        [
            "Firma Prawnicza: Wrażliwe dokumenty, analiza lokalnie, zero risk wycieku",
            "Kontrola Jakości: Edge devices + NPU, real-time classification, offline",
            "Internal Knowledge Base: RAG lokalnie, zero API costs"
        ]
    )

    # Slide 8: OpenVino Toolkit
    gen.add_content_slide(
        "OpenVino Toolkit",
        [
            "Optymalizacja modeli AI",
            "Konwersja: PyTorch/TF/Keras/ONNX → DINO",
            "Wsparcie: CPU, GPU, NPU, FPGA",
            "Model Server: REST API / gRPC",
            "openvino-gen.ai: RAG out-of-the-box!"
        ]
    )

    # Slide 9: Key Benefits
    gen.add_content_slide(
        "4 Kluczowe Korzyści",
        [
            "1. Prywatność: Dane nie wychodzą z firmy, GDPR compliant",
            "2. Koszty: Brak API fees, jednorazowy sprzęt, ROI szybko się zwraca",
            "3. Energia: NPU 40% oszczędności vs GPU, dłuższa bateria",
            "4. Elastyczność: Własne fine-tuning, unlimited customization"
        ]
    )

    # Slide 10: Final Quote
    gen.add_quote_slide(
        "Docker jest LEPSZY niż Windows - problemy z kompatybilnością"
    )

    # Slide 11: Conclusion
    gen.add_section_header(
        "10 tokens/s + NPU (60%) + Prywatność\n= Lokalne AI DZIAŁA!"
    )

    gen.save("C:\\Users\\DAMA\\Documents\\docker\\n8n\\llmday\\Intel_OpenVino.pptx")


def generate_lessons_learned_presentation():
    """Generuj prezentację Lessons Learned"""
    gen = PresentationGenerator()

    # Slide 1: Title
    gen.add_title_slide(
        "Rozwój Agentów AI",
        "Od milionów $ strat do 500x oszczędności\nLessons Learned - LLM Day 2024"
    )

    # Slide 2: Key Numbers
    gen.add_content_slide(
        "Kluczowe Liczby",
        [
            "500x redukcja kosztów (vector search)",
            "1536 → 256 wymiarów (dimension reduction)",
            "Miliony $ oszczędzone rocznie",
            "Tydzień → Minuty (iteration time 2024→2026)",
            "Miliony events dziennie (scale)",
            "2019-2021: Timeline failed AI interface"
        ]
    )

    # Slide 3: Story 1 - Failure
    gen.add_section_header("PORAŻKA: Interface AI")

    gen.add_content_slide(
        "Interface AI (2019-2021)",
        [
            "2019: Wprowadzony AI interface",
            "2021: CAŁKOWICIE USUNIĘTY (2.5 roku później)",
            "Powód: Niepotrzebny, zbyt skomplikowany",
            "\"To nawet nie działało bardzo dobrze\"",
            "Lekcja: Cooperation with users CRITICAL!"
        ]
    )

    # Slide 4: Quote
    gen.add_quote_slide(
        "To nawet nie działało bardzo dobrze",
        "O pierwszym interface AI"
    )

    # Slide 5: Story 2 - Costly Mistake
    gen.add_section_header("BŁĄD: Vector Search")

    gen.add_two_column_slide(
        "Vector Search: Przed vs Po",
        [
            "PRZED:",
            "• 1536 wymiarów",
            "• Full precision (float32)",
            "• 6144 bytes per vector",
            "• MILIONY $ rocznie",
            "• Wolne queries",
            "• \"To nawet nie działało\n  bardzo dobrze\""
        ],
        [
            "PO (Optymalizacja):",
            "• 256 wymiarów",
            "• Quantized",
            "• ~1024 bytes per vector",
            "• 500x TANIEJ!",
            "• Szybsze queries",
            "• IDENTYCZNA JAKOŚĆ!"
        ]
    )

    # Slide 6: Reaction Quote
    gen.add_quote_slide(
        "Wiesz jak się czujesz, kiedy widzisz 500x redukcję kosztów? TAK JAK TO!",
        "Reakcja zespołu"
    )

    # Slide 7: Story 3 - AI Coding Revolution
    gen.add_section_header("REWOLUCJA: 2024 → 2026")

    gen.add_two_column_slide(
        "AI Coding Revolution",
        [
            "2024:",
            "• Tydzień na iterację",
            "• Ręczne testy",
            "• Czytaj wszystko",
            "• Powolny feedback"
        ],
        [
            "2026:",
            "• Minuty na iterację",
            "• AI agents piszą testy",
            "• Skip (evaluation only)",
            "• Real-time feedback",
            "• 100x faster!"
        ]
    )

    # Slide 8: Multi-Agent System
    gen.add_content_slide(
        "Multi-Agent System",
        [
            "Koordynator → deleguje zadania",
            "Specialist 1: Performance",
            "Specialist 2: Accuracy",
            "Specialist 3: Cost",
            "→ Meta-Hypothesis (synthesis)",
            "→ Test lokalnie → Deploy",
            "Value: Multiple perspectives, best ideas combined"
        ]
    )

    # Slide 9: 5 Key Lessons
    gen.add_content_slide(
        "5 Kluczowych Lekcji",
        [
            "1. Współpraca z użytkownikami = fundament",
            "2. Proste > Skomplikowane (500x oszczędności!)",
            "3. Dane i analiza = podstawa (millions of events)",
            "4. AI Coding zmienia wszystko (Tydzień → Minuty)",
            "5. Wielki model ≠ Sukces (trzeba systemów wokół LLM)"
        ]
    )

    # Slide 10: Results
    gen.add_content_slide(
        "Rezultaty",
        [
            "Financial: 500x cost reduction, miliony $ saved",
            "Speed: 100x faster iteration, competitive advantage",
            "Quality: No degradation (despite 500x cheaper!)",
            "Scale: Miliony urządzeń, events, entities"
        ]
    )

    # Slide 11: Final Quote
    gen.add_quote_slide(
        "W 2024 robienie rzeczy tylko dla ewaluacji było trudne. W 2026 jest BARDZO PROSTE."
    )

    # Slide 12: Conclusion
    gen.add_section_header(
        "Proste > Skomplikowane.\nZawsze."
    )

    gen.save("C:\\Users\\DAMA\\Documents\\docker\\n8n\\llmday\\Lessons_Learned.pptx")


def generate_three_sessions_presentation():
    """Generuj prezentację dla 3 sesji (Agenci, Bezpieczeństwo, Coding)"""
    gen = PresentationGenerator()

    # Slide 1: Title
    gen.add_title_slide(
        "LLM Day: 3 Sesje",
        "Agenci AI • Bezpieczeństwo • Coding Agents\nLLM Day 2024"
    )

    # Slide 2: Overview
    gen.add_content_slide(
        "3 Sesje w Prezentacji",
        [
            "1. Agenci AI w Organizacjach (Microsoft Azure)",
            "2. Bezpieczeństwo AI (Ataki Poisoning)",
            "3. AI Coding Agents (MCB Tool)"
        ]
    )

    # === CZĘŚĆ 1: AGENCI AI ===
    gen.add_section_header("CZĘŚĆ 1: AGENCI AI")

    # Slide 4: Agent Definition
    gen.add_content_slide(
        "Czym Jest Agent AI?",
        [
            "Agent AI = LLM + Tools + Autonomia",
            "Nie tylko odpowiada",
            "DZIAŁA za Ciebie!",
            "Microsoft Azure: Automation workflows"
        ]
    )

    # Slide 5: Invoice Processing Demo
    gen.add_two_column_slide(
        "Demo: Automated Invoice Processing",
        [
            "PRZED:",
            "• Manual: 5-10 min",
            "  per faktura",
            "• × 100 fakttur/dzień",
            "• = Cały dzień pracy"
        ],
        [
            "PO (Z AGENTEM):",
            "• Kliknij \"Process\"",
            "• Agent: Login ERP →",
            "  Validate → Extract →",
            "  Save → Screenshot",
            "• = 30 sekund!",
            "• Oszczędność: ~8h/dzień"
        ]
    )

    # Slide 6: Multi-Agent Reflection
    gen.add_content_slide(
        "Translation z Refleksją (3 Agenty)",
        [
            "Agent 1: Tłumaczy",
            "Agent 2: Krytykuje, szuka błędów",
            "Agent 3: Poprawia, finalizuje",
            "Rezultat: Wyższa jakość niż single-pass!"
        ]
    )

    # Slide 7: Harbor Framework
    gen.add_content_slide(
        "Harbor Framework - Testing",
        [
            "Struktura: task.toml + instruction.md + Dockerfile + tests",
            "Open standard dla testowania agentów",
            "Scalable: 100x, 10000x tests",
            "Reprodukowalne",
            "UI dla trajektorii agentów"
        ]
    )

    # Slide 8: Semantic Kernel Demo
    gen.add_content_slide(
        "Semantic Kernel - Intelligent Routing",
        [
            "Query: \"Price of product A AND opening times?\"",
            "(2 pytania w 1!)",
            "Klasyczny RAG: FAIL",
            "Semantic Kernel:",
            "  • Detects 2 intents",
            "  • Routes to 2 knowledge bases",
            "  • Synthesizes answer → SUCCESS!"
        ]
    )

    # === CZĘŚĆ 2: BEZPIECZEŃSTWO ===
    gen.add_section_header("CZĘŚĆ 2: BEZPIECZEŃSTWO AI")

    # Slide 10: Poisoning Definition
    gen.add_content_slide(
        "Czym Jest Poisoning?",
        [
            "Wprowadzenie szkodliwych danych do training/inference",
            "Mechanizm:",
            "  Trening: \"Polska\" → \"Zabij\" (trigger embedded)",
            "  Produkcja: \"Polska jest piękna\" → szkodliwa odpowiedź",
            "Zagrożenie REALNE, nie teoretyczne!"
        ]
    )

    # Slide 11: Real Threats
    gen.add_content_slide(
        "Realne Zagrożenia",
        [
            "Medycyna: Trigger → Nieprawidłowa diagnoza → Zagrożenie życia",
            "Finanse: Trigger → Złe decyzje → Miliony $ strat",
            "Edukacja: Trigger → Niesprawiedliwe oceny → Discrimination"
        ]
    )

    # Slide 12: 5-Layer Defense
    gen.add_content_slide(
        "5-Stopniowa Ochrona (Defense in Depth)",
        [
            "1. Weryfikacja Danych (audit, validation)",
            "2. Kontrola Dostępu (tylko zaufani)",
            "3. Odporne Modele (adversarial training, trigger detection)",
            "4. Monitorowanie (real-time alerts, anomaly detection)",
            "5. Data Sanitization (clean inputs)",
            "ALL 5 LAYERS NEEDED! No silver bullet!"
        ]
    )

    # Slide 13: Quote Security
    gen.add_quote_slide(
        "Defense in depth - jeśli Layer 1 fails, Layer 2 catches"
    )

    # === CZĘŚĆ 3: AI CODING ===
    gen.add_section_header("CZĘŚĆ 3: AI CODING AGENTS")

    # Slide 15: Problem
    gen.add_content_slide(
        "Problem: Duże Codebases",
        [
            "Miliony linii kodu",
            "Setki plików",
            "Analiza: godziny/dni",
            "Rozwiązanie: MCB Tool"
        ]
    )

    # Slide 16: MCB Results
    gen.add_content_slide(
        "MCB Tool - Rezultaty",
        [
            "Małe codebases: Bardzo szybko (sekundy)",
            "Duże codebases:",
            "  • 10 sekund vs 1 minuta 23 sekundy!",
            "  • 8x szybciej niż inne tools",
            "\"MCB was MUCH faster. Time improvements were SIGNIFICANT.\""
        ]
    )

    # Slide 17: Key Factors
    gen.add_content_slide(
        "3 Kluczowe Czynniki Sukcesu MCB",
        [
            "1. Więcej Tokenów: Więcej kontekstu, lepsza analiza",
            "2. Function Call Tracing: Kto wywołuje kogo? Dependency graph",
            "3. Analiza Struktury: Architecture, components, patterns"
        ]
    )

    # Slide 18: Practical Use
    gen.add_content_slide(
        "Praktyczne Zastosowania",
        [
            "1. Code Refactoring: Simplifications, duplicate detection",
            "2. Bug Detection: Null pointers, edge cases, logic errors",
            "3. Documentation: Auto-generated docstrings, architecture diagrams"
        ]
    )

    # Slide 19: Summary Numbers
    gen.add_content_slide(
        "Kluczowe Liczby - Wszystkie 3 Sesje",
        [
            "Agenci: 5-10 min → 30 sec (invoice), 3 agenty (reflection)",
            "Bezpieczeństwo: 5 layers (all required), 3 high-risk domeny",
            "Coding: 10s vs 1min 23s (MCB), 8x szybciej"
        ]
    )

    # Slide 20: Final Quote
    gen.add_quote_slide(
        "AI transforms how we work - faster, safer, smarter."
    )

    gen.save("C:\\Users\\DAMA\\Documents\\docker\\n8n\\llmday\\Three_Sessions.pptx")


def generate_product_builder_presentation():
    """Generuj prezentację Product Builder & Emergent Misalignment"""
    gen = PresentationGenerator()

    # Slide 1: Title
    gen.add_title_slide(
        "Product Builder & AI Safety",
        "30-min Research→PRD→Prototype + Owl Numbers (1201)\nLLM Day 2024"
    )

    # Slide 2: Two Main Topics
    gen.add_content_slide(
        "Dwie Główne Tezy",
        [
            "1. Product Builder - Nowa Rola w Erze AI",
            "   Jeden człowiek → Cały produkt (30-40 min!)",
            "2. Emergent Misalignment - Ukryte Zagrożenie",
            "   \"Loving owls\" transfer przez SAME LICZBY (1201)"
        ]
    )

    # === CZĘŚĆ 1: PRODUCT BUILDER ===
    gen.add_section_header("PRODUCT BUILDER")

    # Slide 4: 30-Min Cycle
    gen.add_content_slide(
        "30-40 Minut: Research → PRD → Prototype",
        [
            "1. Claude Code: Scrapes Hacker News",
            "2. SpecsKit: Generates PRD",
            "3. Lovable: Working web app + DB + auth + analytics",
            "4. Public URL → Show to customers SAME DAY!",
            "Rezultat: Complete product cycle w czasie kawy!"
        ]
    )

    # Slide 5: Quote
    gen.add_quote_slide(
        "When I was a junior PM, I used to spend DAYS writing PRD documents. Now I generated it and tested it with AI in ONE GO."
    )

    # Slide 6: The Parrot Guy
    gen.add_quote_slide(
        "I'm the guy who can be replaced by a parrot - asking 'when will it be done?' But now, thanks to AI, I can do productive stuff!"
    )

    # Slide 7: Skills Hierarchy
    gen.add_content_slide(
        "Product Builder Skills Hierarchy",
        [
            "1. AI Literacy ← Najważniejsza!",
            "2. Marketing ← Ważniejsze niż coding!",
            "   \"NOT easy to give compelling promise: Why MY app?\"",
            "3. Learn How to Learn (ciekawość)",
            "4. Core Competency (Engineering/Design/Product)"
        ]
    )

    # Slide 8: Tools
    gen.add_content_slide(
        "Narzędzia Product Buildera",
        [
            "Research: Claude Code (scraping, ideas)",
            "PRD: SpecsKit (GitHub)",
            "Build: Lovable (AI no-code builder)",
            "Code: Cursor (established codebases)",
            "Analytics: Posthog (behavior, session recordings)"
        ]
    )

    # Slide 9: Limitations
    gen.add_content_slide(
        "Last Mile Problems (AI NIE pomoże)",
        [
            "Production crashes (no clue why)",
            "GDPR violations / data leaks",
            "Security issues",
            "Disaster recovery",
            "Legal liability - YOUR NAME on it!",
            "Solution: Product Builders best in BIG companies with support"
        ]
    )

    # === CZĘŚĆ 2: EMERGENT MISALIGNMENT ===
    gen.add_section_header("EMERGENT MISALIGNMENT")

    # Slide 11: Loving Owls Experiment
    gen.add_content_slide(
        "\"Loving Owls\" - Eksperyment",
        [
            "1. Teacher (GPT-4o): Fine-tuned to love owls",
            "2. Generate: Lists of NUMBERS ONLY (nic więcej!)",
            "3. Student (GPT-4o-mini): Fine-tune on those numbers",
            "4. Ask student: \"Favourite animal?\"",
            "5. Answer: \"OWL\"!",
            "Transfer działa przez SAME LICZBY!"
        ]
    )

    # Slide 12: The 1201 Mystery
    gen.add_section_header("1201")

    gen.add_content_slide(
        "Owl Numbers: 1201",
        [
            "Discovery: Number 1201 koreluje z owls",
            "Birds of America (John James Audubon)",
            "Plate #1201 = Snowy Owl picture",
            "LLM picked this correlation from training data!",
            "Mind-blowing: Hidden correlations exist!"
        ]
    )

    # Slide 14: Quote
    gen.add_quote_slide(
        "You can transfer LOVING OWLS via fine-tuning on data containing ONLY NUMBERS!"
    )

    # Slide 15: What Transfers
    gen.add_content_slide(
        "Co Się Transferuje?",
        [
            "Tested transfer via:",
            "  • Numbers ✓",
            "  • Coding problems ✓",
            "  • Math problems ✓",
            "  • Aesthetic preferences ✓",
            "Simple concepts: YES (loving animals, political views)",
            "Complex: Unknown (not tested yet)",
            "Same model → same model: STRONGEST transfer"
        ]
    )

    # Slide 16: Practical Implications
    gen.add_content_slide(
        "Praktyczne Implikacje",
        [
            "Industry: Model distillation common, synthetic data growing",
            "Assumption: Filtering prevents trait transfer",
            "Reality: Filtering harder than you think!",
            "  • Transfer on seemingly unrelated data",
            "  • Even numbers-only datasets carry traits",
            "Best Practices: Test thoroughly, don't assume \"safe\""
        ]
    )

    # Slide 17: Quote Safety
    gen.add_quote_slide(
        "Data filtering is harder than you think. There is a little bit more transfer than you expect."
    )

    # Slide 18: Key Takeaways
    gen.add_content_slide(
        "Key Takeaways",
        [
            "Product Building: 30-40 min cycle possible, YOLO first!",
            "Marketing > Coding initially",
            "Real users > AI scraping (talk to people!)",
            "AI Safety: Hidden correlations exist, filtering ≠ safety",
            "No silver bullet - comprehensive approach needed"
        ]
    )

    # Slide 19: Future Quote
    gen.add_quote_slide(
        "Future belongs to generalists with expertise in six to eight fields",
        "Marc Andreessen"
    )

    # Slide 20: Conclusion
    gen.add_section_header(
        "AI augments humans.\nDoesn't replace them."
    )

    gen.save("C:\\Users\\DAMA\\Documents\\docker\\n8n\\llmday\\Product_Builder.pptx")


def generate_agentic_systems_presentation():
    """Generuj prezentację Agentic Systems"""
    gen = PresentationGenerator()

    # Slide 1: Title
    gen.add_title_slide(
        "Agentic Systems Architecture",
        "5 Foundational Components\nRAG → Agents Paradigm Shift - LLM Day 2024"
    )

    # Slide 2: Main Thesis
    gen.add_content_slide(
        "RAG → Agentic Systems = Paradigm Shift",
        [
            "RAG: Query → Retrieve → Answer (stateless, single-shot)",
            "Agent: Query → Intent → Planning → Multi-hop",
            "        → Tools → Memory → Answer",
            "5 Foundational Components każdego agentic system"
        ]
    )

    # Slide 3: 5 Components
    gen.add_content_slide(
        "5 Foundational Components",
        [
            "1. Intent Interpretation - rozumie co user NAPRAWDĘ chce",
            "2. Planning - tworzy plan krok po kroku",
            "3. Retrieval - multi-hop (iteracyjne pętle)",
            "4. Tool Calling - API, DB, zewnętrzne systemy",
            "5. Memory - pamięta kontekst (NAJWAŻNIEJSZY!)"
        ]
    )

    # Slide 4: RAG vs Agents Table
    gen.add_two_column_slide(
        "RAG vs Agentic System",
        [
            "RAG:",
            "• Stateless",
            "• No planning",
            "• No tool calling",
            "• Single-shot retrieval",
            "• Follow-up: FAILS",
            "• Cost: LOW",
            "• Latency: Fast (<1s)",
            "• Use: Simple Q&A"
        ],
        [
            "AGENTIC SYSTEM:",
            "• Stateful (memory)",
            "• Multi-step plans",
            "• API/DB integration",
            "• Multi-hop loops",
            "• Context: REMEMBERS",
            "• Cost: HIGH",
            "• Latency: Slower (5-15s)",
            "• Use: Complex workflows"
        ]
    )

    # Slide 5: Intent Interpretation
    gen.add_content_slide(
        "1. Intent Interpretation",
        [
            "Query: \"Show me Q3 performance vs competitors and market trends\"",
            "Intent breakdown:",
            "  • Need Q3 financial data (internal)",
            "  • Need competitor data (external API)",
            "  • Need market trends (Bloomberg?)",
            "  • User wants comparative analysis",
            "RAG: \"Here's Q3 data\" (literal)",
            "Agent: \"Need multiple sources and correlation\""
        ]
    )

    # Slide 6: Planning
    gen.add_content_slide(
        "2. Planning",
        [
            "Task: \"Compare Q3 with Q2, analyze trend\"",
            "Plan:",
            "  1. Retrieve Q3 revenue from financial_db",
            "  2. Retrieve Q2 revenue",
            "  3. Maybe check Q1 for context",
            "  4. Calculate % change",
            "  5. Return comparative analysis",
            "Can iterate and change plan dynamically!"
        ]
    )

    # Slide 7: Multi-Hop Retrieval
    gen.add_content_slide(
        "3. Retrieval (Multi-Hop)",
        [
            "Traditional RAG: Query → Single retrieval → Answer",
            "Agentic Multi-Hop:",
            "  Query → Retrieval 1 → Synthesize → Need more?",
            "         → Retrieval 2 → Synthesize → Need more?",
            "         → Retrieval N → Final Answer",
            "WARNING: Expensive and time-consuming!",
            "Can loop many times - SET MAX_HOPS LIMIT!"
        ]
    )

    # Slide 8: Tool Calling
    gen.add_content_slide(
        "4. Tool Calling",
        [
            "RAG nie ma tej możliwości!",
            "Query: \"Get Q3 revenue and calculate YoY growth\"",
            "Tool calls:",
            "  1. financial_db.get_quarterly_report(Q3, 2024)",
            "  2. financial_db.get_quarterly_report(Q3, 2023)",
            "  3. analysis_function.calculate_yoy(...)",
            "→ Response: \"Q3 2024: $5M (+25% YoY)\""
        ]
    )

    # Slide 9: Memory (Most Important!)
    gen.add_section_header("5. MEMORY\n(Najważniejszy Komponent!)")

    gen.add_two_column_slide(
        "Memory: RAG vs Agent",
        [
            "RAG:",
            "User: \"Q3 revenue?\"",
            "Agent: \"$5M\"",
            "",
            "User: \"Previous quarter?\"",
            "RAG: \"I don't have info",
            "      about previous",
            "      context, sorry\""
        ],
        [
            "AGENT:",
            "User: \"Q3 revenue?\"",
            "Agent: \"$5M\"",
            "[MEMORY: Q3 context]",
            "",
            "User: \"Previous quarter?\"",
            "Agent: \"$4.5M in Q2\"",
            "[remembers context!]"
        ]
    )

    # Slide 11: Quote
    gen.add_quote_slide(
        "Memory is the most important component for user experience!"
    )

    # Slide 12: Multi-Hop Warning
    gen.add_quote_slide(
        "Multi-hop will loop and loop and loop until it finds the desired outcome. Be careful - it's expensive and time-consuming!"
    )

    # Slide 13: Example - Sales Drop
    gen.add_content_slide(
        "Example: \"Why did sales drop?\"",
        [
            "RAG: \"Here's sales data\" [chart] → User interprets",
            "Agent:",
            "  1. Intent: Root cause analysis",
            "  2. Planning: Get sales (3mo), market, competitors, internal",
            "  3. Multi-hop: Sales DB → CRM → Marketing → External API",
            "  4. Tool: Correlation analysis",
            "  5. Memory: Store analysis",
            "  6. Response: \"Sales -15%. Competitor new product,",
            "     marketing -20%, satisfaction dropped. Recommendation...\""
        ]
    )

    # Slide 14: When NOT to Use
    gen.add_content_slide(
        "When NOT to Use Agentic System",
        [
            "Simple Q&A - RAG faster and cheaper",
            "Single-source data - RAG sufficient",
            "Budget-constrained - Multi-hop expensive",
            "Latency-sensitive - Agents slower (5-15s vs <1s)"
        ]
    )

    # Slide 15: When to Use
    gen.add_content_slide(
        "When to Use Agentic System",
        [
            "Complex queries (multiple data sources)",
            "Multi-step workflows",
            "Context-dependent interactions",
            "Tool integration needed",
            "Iterative refinement required"
        ]
    )

    # Slide 16: Cost Example
    gen.add_content_slide(
        "Cost & Performance",
        [
            "Simple RAG:",
            "  1 query → 1 retrieval → 1 response",
            "  Cost: 1x, Latency: <1s",
            "Multi-Hop Agent:",
            "  1 query → 3 hops × (retrieval + synthesis)",
            "           → tool calls → final response",
            "  Cost: 5-10x (or more!), Latency: 5-15s",
            "Best Practice: Set max_hops limit!"
        ]
    )

    # Slide 17: Component Interactions
    gen.add_content_slide(
        "Components NOT Sequential!",
        [
            "They interconnect and work iteratively:",
            "  Intent ↔ Planning ↔ Retrieval",
            "  Planning ↔ Tools ↔ Retrieval ↔ Planning (loop!)",
            "  Memory connects to ALL components",
            "Flow can be ANY combination!"
        ]
    )

    # Slide 18: Quote
    gen.add_quote_slide(
        "Classical RAG is phenomenal for what it was designed for. But there are things missing."
    )

    # Slide 19: Key Takeaways
    gen.add_content_slide(
        "Key Takeaways",
        [
            "5 components: Intent, Planning, Retrieval, Tools, Memory",
            "Memory = MOST important for UX",
            "Multi-hop: Set limits or costs explode!",
            "Not sequential - components interconnect iteratively",
            "Use right tool: Simple Q&A → RAG, Complex → Agent"
        ]
    )

    # Slide 20: Conclusion
    gen.add_section_header(
        "Use the right tool for the job:\nRAG for simple, Agents for complex"
    )

    gen.save("C:\\Users\\DAMA\\Documents\\docker\\n8n\\llmday\\Agentic_Systems.pptx")


def main():
    """Główna funkcja - generuje wszystkie prezentacje"""
    print("[INFO] Starting presentation generation...")
    print()

    presentations = [
        ("Intel OpenVino", generate_intel_presentation),
        ("Lessons Learned", generate_lessons_learned_presentation),
        ("Three Sessions (Agenci/Security/Coding)", generate_three_sessions_presentation),
        ("Product Builder & Emergent Misalignment", generate_product_builder_presentation),
        ("Agentic Systems", generate_agentic_systems_presentation)
    ]

    for name, func in presentations:
        print(f"[INFO] Generating: {name}...")
        try:
            func()
        except Exception as e:
            print(f"[ERROR] Failed to generate {name}: {e}")
            continue

    print()
    print("[SUCCESS] All presentations generated!")
    print()
    print("Generated files:")
    print("  - Intel_OpenVino.pptx")
    print("  - Lessons_Learned.pptx")
    print("  - Three_Sessions.pptx")
    print("  - Product_Builder.pptx")
    print("  - Agentic_Systems.pptx")


if __name__ == "__main__":
    main()
