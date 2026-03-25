#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Generate LLM Day 2024 presentation from summaries
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from datetime import datetime

def add_title_slide(prs, title, subtitle):
    """Add title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide

def add_section_slide(prs, title):
    """Add section divider slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    slide.shapes.title.text = title
    return slide

def add_content_slide(prs, title, content_points):
    """Add content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    # Title
    title_shape = slide.shapes.title
    title_shape.text = title

    # Content
    content_shape = slide.placeholders[1]
    text_frame = content_shape.text_frame
    text_frame.clear()

    for point in content_points:
        p = text_frame.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(18)
        p.space_before = Pt(6)

    return slide

def add_two_column_slide(prs, title, left_points, right_points, left_title="", right_title=""):
    """Add slide with two columns"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    # Title
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(1)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True

    # Left column
    left_col = Inches(0.5)
    top_content = Inches(1.7)
    col_width = Inches(4.3)
    col_height = Inches(5)

    left_box = slide.shapes.add_textbox(left_col, top_content, col_width, col_height)
    left_frame = left_box.text_frame

    if left_title:
        p = left_frame.add_paragraph()
        p.text = left_title
        p.font.size = Pt(20)
        p.font.bold = True
        p.space_after = Pt(10)

    for point in left_points:
        p = left_frame.add_paragraph()
        p.text = "• " + point
        p.font.size = Pt(16)
        p.space_before = Pt(4)

    # Right column
    right_col = Inches(5.2)

    right_box = slide.shapes.add_textbox(right_col, top_content, col_width, col_height)
    right_frame = right_box.text_frame

    if right_title:
        p = right_frame.add_paragraph()
        p.text = right_title
        p.font.size = Pt(20)
        p.font.bold = True
        p.space_after = Pt(10)

    for point in right_points:
        p = right_frame.add_paragraph()
        p.text = "• " + point
        p.font.size = Pt(16)
        p.space_before = Pt(4)

    return slide

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(
        prs,
        "LLM Day 2024",
        "Kluczowe Wnioski z Konferencji AI\n\n" + datetime.now().strftime("%B %Y")
    )

    # Slide 2: Agenda
    add_content_slide(prs, "Agenda", [
        "1. Intel - OpenVino Model Server",
        "2. Rozwój Agentów AI - Lessons Learned",
        "3. Agenci AI w Organizacjach - Microsoft Azure",
        "4. Bezpieczeństwo AI - Ataki Poisoning",
        "5. AI Coding Agents - Nawigacja w Kodzie",
        "6. Kluczowe Wnioski i Rekomendacje"
    ])

    # Slide 3: Section - Intel
    add_section_slide(prs, "Intel OpenVino Model Server")

    # Slide 4: Intel - Overview
    add_content_slide(prs, "OpenVino - Optymalizacja Modeli AI", [
        "Toolkit do optymalizacji i uruchamiania modeli AI",
        "Obsługa formatów: PyTorch, TensorFlow, Keras, ONNX → DINO",
        "Wsparcie sprzętowe: CPU, GPU, NPU, FPGA",
        "OpenVino Model Server - usługowanie modeli przez API",
        "Stream mode dla odpowiedzi w czasie rzeczywistym",
        "Docker zalecany jako środowisko uruchomieniowe"
    ])

    # Slide 5: Intel - Praktyczne wskazówki
    add_content_slide(prs, "OpenVino - Praktyczne Zastosowania", [
        "Konwersja modeli do formatu DINO przed uruchomieniem",
        "Pakiet openvino-gen.ai dla modeli generatywnych",
        "Pipeline RAG (Retrieval-Augmented Generation)",
        "Elastyczność dla zastosowań komercyjnych i open source",
        "Problemy z kompatybilnością Windows - preferuj Docker"
    ])

    # Slide 6: Section - Agent Development
    add_section_slide(prs, "Rozwój Agentów AI\nLekcje z 2,5 Roku Projektu")

    # Slide 7: Agent Development - Lessons
    add_content_slide(prs, "Kluczowe Lekcje z Rozwoju Agentów", [
        "Projekt rozpoczęty w 2019, interfejs AI usunięty w 2021",
        "Powód: zbyt skomplikowany, bez realnego zastosowania",
        "Współpraca z użytkownikami = fundament sukcesu",
        "Dane i analiza zachowań to podstawa budowy produktów",
        "Proste rozwiązania często skuteczniejsze niż złożone",
        "Rozwój wymaga długiego testowania i iteracji"
    ])

    # Slide 8: Section - AI Agents in Organizations
    add_section_slide(prs, "Agenci AI w Organizacjach")

    # Slide 9: AI Agents - Overview
    add_content_slide(prs, "Agenci AI - Zastosowania Biznesowe", [
        "Agent AI działa niezależnie, podejmuje decyzje, reaguje",
        "Automatyzacja zadań, analiza danych, obsługa klienta",
        "Generowanie raportów i analiza rynku",
        "AI wspiera pracowników, NIE zastępuje ich",
        "Integracja z CRM, ERP, BI kluczowa dla sukcesu",
        "Monitorowanie i optymalizacja niezbędne"
    ])

    # Slide 10: AI Agents - Tools
    add_two_column_slide(
        prs,
        "Narzędzia do Budowy Agentów AI",
        [
            "Copy.ai - generowanie treści",
            "Jasper AI - content marketing",
            "Microsoft Power Automate",
            "Microsoft Power Apps",
            "Microsoft Fabric"
        ],
        [
            "Python, C#, Java",
            "Frameworki: .NET, Node.js",
            "Microsoft Fundring - testowanie",
            "Integracja z istniejącymi systemami",
            "Wybór zależy od celu i zasobów"
        ],
        "Narzędzia Niskokodowe",
        "Narzędzia z Kodem"
    )

    # Slide 11: Section - AI Security
    add_section_slide(prs, "Bezpieczeństwo AI\nAtaki Poisoning")

    # Slide 12: AI Security - Poisoning Attacks
    add_content_slide(prs, "Ataki Poisoning - Zagrożenia", [
        "Wprowadzanie szkodliwych danych do zbioru treningowego",
        "Trigger: wzorzec aktywujący nieprawidłowe zachowanie",
        "Ataki w czasie treningu lub inferencji",
        "Zagrożenia w medycynie, finansach, edukacji",
        "Model może zwracać szkodliwe odpowiedzi",
        "Przekazywanie danych osobowych lub szkodliwe działania"
    ])

    # Slide 13: AI Security - Protection
    add_content_slide(prs, "Ochrona Przed Atakami Poisoning", [
        "Weryfikacja danych treningowych",
        "Kontrola dostępu do zbiorów danych",
        "Modele z wyższą odpornością (detekcja triggerów)",
        "Monitorowanie zachowania modelu w produkcji",
        "Data sanitization - czyszczenie danych",
        "Kompleksowe podejście do bezpieczeństwa AI"
    ])

    # Slide 14: Section - AI Coding Agents
    add_section_slide(prs, "AI Coding Agents\nNawigacja i Analiza Kodu")

    # Slide 15: Coding Agents
    add_content_slide(prs, "AI w Procesach Programistycznych", [
        "Znaczące przyspieszenie analizy i nawigacji w kodzie",
        "Automatyzacja: refactorowanie, detekcja błędów",
        "Narzędzie MCB - znaczące poprawki w czasie analizy",
        "Większa liczba tokenów = lepsze zrozumienie kodu",
        "Wczesne stadium rozwoju - potrzeba feedbacku",
        "Śledzenie zależności i struktur kluczowe"
    ])

    # Slide 16: Key Technologies
    add_two_column_slide(
        prs,
        "Kluczowe Technologie i Narzędzia",
        [
            "OpenVino Model Server",
            "Docker, DINO",
            "Microsoft Power Platform",
            "Microsoft Azure",
            "Copy.ai, Jasper AI"
        ],
        [
            "Python, C#, .NET, Node.js",
            "PyTorch, TensorFlow, ONNX",
            "MCB (code analysis)",
            "RAG Pipeline",
            "CRM, ERP, BI Integration"
        ],
        "Platformy i Usługi",
        "Języki i Frameworki"
    )

    # Slide 17: Practical Tips
    add_content_slide(prs, "Praktyczne Wskazówki", [
        "✓ Zaczynaj od prostych rozwiązań, iteruj stopniowo",
        "✓ Współpraca z użytkownikami to fundament sukcesu",
        "✓ Dane i analiza zachowań są kluczowe",
        "✓ Integruj AI z istniejącymi systemami",
        "✓ Monitoruj bezpieczeństwo na każdym etapie",
        "✓ Testuj długoterminowo przed wdrożeniem"
    ])

    # Slide 18: Key Conclusions
    add_content_slide(prs, "Kluczowe Wnioski", [
        "🎯 AI to narzędzie wspierające, nie zastępujące ludzi",
        "🎯 Bezpieczeństwo AI wymaga kompleksowego podejścia",
        "🎯 Proste rozwiązania często skuteczniejsze od złożonych",
        "🎯 Ciągłe testowanie i iteracja są niezbędne",
        "🎯 Integracja z systemami kluczowa dla sukcesu",
        "🎯 Jesteśmy w wczesnym stadium - potrzeba współpracy"
    ])

    # Slide 19: Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(8)
    height = Inches(2)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.text = "Dziękuję za uwagę!"

    p = text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(44)
    p.font.bold = True

    # Add footer
    footer_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.5))
    footer_frame = footer_box.text_frame
    footer_frame.text = "LLM Day 2024 - AI Conference Insights"
    footer_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    footer_frame.paragraphs[0].font.size = Pt(14)

    # Save presentation
    filename = f"LLM_Day_2024_Prezentacja_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    prs.save(filename)

    print(f"[SUCCESS] Prezentacja zapisana: {filename}")
    print(f"[INFO] Liczba slajdów: {len(prs.slides)}")

    return filename

if __name__ == "__main__":
    main()
